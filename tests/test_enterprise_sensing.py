from __future__ import annotations

import pytest
from pptx import Presentation
from io import BytesIO
from streamlit.testing.v1 import AppTest

from src.models.enterprise import (
    EnterpriseDataDimension,
    EnterpriseEvidenceCategory,
    EnterpriseEvidenceItem,
    EnterpriseReviewStatus,
    EnterpriseSensitivity,
    EnterpriseStatementType,
)
from src.services.enterprise_sensing import (
    MAX_UPLOAD_BYTES,
    EnterpriseSensingError,
    accepted_unredacted_entries,
    company_strategy_gate_reasons,
    confirm_enterprise_artifact,
    diagnosis_title_from_symptoms,
    delete_enterprise_entry,
    enterprise_item_from_upload,
    extract_document_text,
    new_enterprise_artifact,
    review_enterprise_entry,
    upsert_enterprise_entry,
)
from src.state.project import ProjectState


def make_strategy_project() -> ProjectState:
    return ProjectState(
        project_name="Molecular Diagnostics Strategy",
        industry="Molecular Diagnostics",
        region="China",
        target_company="Demo Diagnostics",
        company_strategy_enabled=True,
        company_strategy_objective="Protect PCR cash flow while testing a digital PCR expansion",
        research_objective="Assess competitors, drivers, and future trends",
        time_horizon="2026-2030",
    )


def make_item(
    sensitivity: EnterpriseSensitivity = EnterpriseSensitivity.REDACTED_DEMO,
) -> EnterpriseEvidenceItem:
    return EnterpriseEvidenceItem(
        title="Channel observation",
        category=EnterpriseEvidenceCategory.SALES_CHANNEL,
        statement_type=EnterpriseStatementType.OBSERVATION,
        content="Three simulated distributors report demand for simpler workflows.",
        source_owner="Demo sales lead",
        strategic_relevance="Tests whether workflow integration supports the expansion goal.",
        sensitivity=sensitivity,
    )


def accepted_artifact(project: ProjectState, *, sensitivity=EnterpriseSensitivity.REDACTED_DEMO):
    artifact = upsert_enterprise_entry(None, project, make_item(sensitivity))
    item_id = artifact.entries[0].enterprise_evidence_id
    return review_enterprise_entry(
        artifact,
        item_id,
        EnterpriseReviewStatus.ACCEPTED,
        "Suitable for the demo",
    )


def test_enterprise_artifact_requires_an_accepted_item() -> None:
    project = make_strategy_project()
    artifact = new_enterprise_artifact(project)

    with pytest.raises(EnterpriseSensingError, match="至少需要一条"):
        confirm_enterprise_artifact(
            artifact,
            project,
            consent_to_model_processing=True,
            public_demo_acknowledged=True,
        )


def test_confirmed_enterprise_artifact_unlocks_strategy_eligibility() -> None:
    project = make_strategy_project()
    artifact = accepted_artifact(project)
    confirmed = confirm_enterprise_artifact(
        artifact,
        project,
        consent_to_model_processing=True,
        public_demo_acknowledged=True,
    )
    project = project.model_copy(update={"enterprise_sensing_artifact": confirmed})

    assert confirmed.human_confirmed is True
    assert company_strategy_gate_reasons(project) == []


def test_strategy_change_invalidates_enterprise_snapshot() -> None:
    project = make_strategy_project()
    confirmed = confirm_enterprise_artifact(
        accepted_artifact(project),
        project,
        consent_to_model_processing=True,
        public_demo_acknowledged=True,
    )
    project = project.model_copy(
        update={
            "enterprise_sensing_artifact": confirmed,
            "company_strategy_objective": "Prioritize NGS services instead",
        }
    )

    assert any("战略目标已变化" in reason for reason in company_strategy_gate_reasons(project))


def test_public_demo_rejects_confidential_enterprise_input() -> None:
    project = make_strategy_project()
    artifact = accepted_artifact(project, sensitivity=EnterpriseSensitivity.CONFIDENTIAL)

    with pytest.raises(EnterpriseSensingError, match="仍标记"):
        confirm_enterprise_artifact(
            artifact,
            project,
            consent_to_model_processing=True,
            public_demo_acknowledged=True,
        )


def test_public_demo_also_rejects_unredacted_internal_input() -> None:
    project = make_strategy_project()
    artifact = accepted_artifact(project, sensitivity=EnterpriseSensitivity.INTERNAL)

    with pytest.raises(EnterpriseSensingError, match="仍标记"):
        confirm_enterprise_artifact(
            artifact,
            project,
            consent_to_model_processing=True,
            public_demo_acknowledged=True,
        )


def test_rejected_confidential_input_does_not_block_public_demo() -> None:
    project = make_strategy_project()
    artifact = accepted_artifact(project)
    confidential = make_item(EnterpriseSensitivity.CONFIDENTIAL).model_copy(
        update={"title": "Rejected confidential file"}
    )
    artifact = upsert_enterprise_entry(artifact, project, confidential)
    artifact = review_enterprise_entry(
        artifact,
        confidential.enterprise_evidence_id,
        EnterpriseReviewStatus.REJECTED,
        "Not suitable for public demo",
    )

    assert accepted_unredacted_entries(artifact) == []
    confirmed = confirm_enterprise_artifact(
        artifact,
        project,
        consent_to_model_processing=True,
        public_demo_acknowledged=True,
    )
    assert confirmed.human_confirmed is True


def test_enterprise_entry_can_be_deleted() -> None:
    project = make_strategy_project()
    artifact = upsert_enterprise_entry(None, project, make_item())
    item_id = artifact.entries[0].enterprise_evidence_id

    cleaned = delete_enterprise_entry(artifact, item_id)

    assert cleaned.entries == []
    assert cleaned.human_confirmed is False


def test_text_and_csv_uploads_are_extracted_without_model_calls() -> None:
    assert extract_document_text("note.txt", "text/plain", "客户反馈".encode()) == "客户反馈"
    assert extract_document_text("data.csv", "text/csv", b"segment,demand\nhospital,high") == (
        "segment | demand\nhospital | high"
    )

    item = enterprise_item_from_upload(
        file_name="note.txt",
        mime_type="text/plain",
        data=b"redacted enterprise observation",
        source_owner="Product lead",
        strategic_relevance="Tests product-market fit",
        sensitivity=EnterpriseSensitivity.REDACTED_DEMO,
        data_dimension=EnterpriseDataDimension.CUSTOMER_PENETRATION,
        reporting_period="2026Q1",
    )

    assert item.input_method == "file"
    assert item.file_sha256
    assert item.data_dimension == EnterpriseDataDimension.CUSTOMER_PENETRATION
    assert item.reporting_period == "2026Q1"
    assert item.review_status == EnterpriseReviewStatus.NEEDS_REVIEW


def test_enterprise_upload_limit_is_300_mb() -> None:
    assert MAX_UPLOAD_BYTES == 300 * 1024 * 1024


def test_diagnosis_title_is_derived_from_current_symptoms() -> None:
    symptoms = "重点医院客户渗透率偏低。\n现有渠道转化效率未达到计划。"

    assert diagnosis_title_from_symptoms(symptoms) == (
        "重点医院客户渗透率偏低。 现有渠道转化效率未达到计划。"
    )
    assert diagnosis_title_from_symptoms("   ") == ""


def test_oversized_enterprise_upload_is_rejected(monkeypatch) -> None:
    monkeypatch.setattr("src.services.enterprise_sensing.MAX_UPLOAD_BYTES", 4)

    with pytest.raises(EnterpriseSensingError, match="300MB"):
        enterprise_item_from_upload(
            file_name="orders.csv",
            mime_type="text/csv",
            data=b"12345",
            source_owner="Commercial operations",
            strategic_relevance="Tests channel performance",
            sensitivity=EnterpriseSensitivity.REDACTED_DEMO,
        )


def test_powerpoint_upload_is_extracted_without_model_calls() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "企业战略输入"
    slide.placeholders[1].text = "渠道反馈与客户需求"
    buffer = BytesIO()
    presentation.save(buffer)

    text = extract_document_text(
        "strategy.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        buffer.getvalue(),
    )

    assert "企业战略输入" in text
    assert "渠道反馈与客户需求" in text


def test_enterprise_sensing_page_renders_for_case_project() -> None:
    app = AppTest.from_file("app.py").run(timeout=10)
    next(button for button in app.button if button.label == "加载案例展示").click().run(
        timeout=10
    )
    app.selectbox[0].set_value("enterprise_sensing").run(timeout=10)

    assert not app.exception
    assert [item.value for item in app.subheader] == [
        "A. 企业自我诊断问题",
        "B. 分层上传脱敏企业文件",
        "C. 人工审核企业资料",
        "D. 权限确认与战略路径资格",
    ]
    labels = [item.label for item in app.text_area]
    assert "当前表现、症状或已观察到的问题" in labels
    assert not any(item.label == "最需要改进或验证的方面" for item in app.text_input)
    assert any("单文件不超过300MB" in item.value for item in app.caption)
